#!/usr/bin/env python3
"""Generate db_construct_slides.pptx from slide content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette (matches HTML dark theme) ──────────────────────────────────────
BG       = RGBColor(0x0d, 0x11, 0x17)   # slide background
CARD     = RGBColor(0x16, 0x1b, 0x25)   # card/box background
ACCENT   = RGBColor(0x38, 0xbd, 0xf8)   # blue accent (headings)
ACCENT2  = RGBColor(0x34, 0xd3, 0x99)   # green
ACCENT3  = RGBColor(0xfb, 0xbf, 0x24)   # amber
ACCENT4  = RGBColor(0xf8, 0x71, 0x71)   # red
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)
CODE_BG  = RGBColor(0x0a, 0x0e, 0x14)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide  = prs.slides.add_slide(layout)
    fill   = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_textbox(slide, left, top, width, height, text,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True,
                italic=False, bg_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txBox

def add_label(slide, text, top=Inches(0.18)):
    add_textbox(slide, Inches(0.4), top, Inches(4), Inches(0.35),
                text, font_size=10, color=ACCENT, bold=True)

def add_title(slide, text, top=Inches(0.45)):
    add_textbox(slide, Inches(0.4), top, Inches(12.5), Inches(0.65),
                text, font_size=26, bold=True, color=WHITE)

def add_divider(slide, top):
    from pptx.util import Pt as PtU
    line = slide.shapes.add_connector(1, Inches(0.4), top, Inches(12.9), top)
    line.line.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    line.line.width = Pt(0.75)

def add_card(slide, left, top, width, height, title, body,
             title_color=ACCENT2, font_body=12):
    # card background
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    box.line.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    box.line.width = Pt(0.75)

    # title
    add_textbox(slide, left + Inches(0.12), top + Inches(0.08),
                width - Inches(0.2), Inches(0.35),
                title, font_size=13, bold=True, color=title_color)
    # body
    add_textbox(slide, left + Inches(0.12), top + Inches(0.42),
                width - Inches(0.2), height - Inches(0.5),
                body, font_size=font_body, color=WHITE, word_wrap=True)

def add_code_box(slide, left, top, width, height, code_text, font_size=9):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.08),
                width - Inches(0.15), height - Inches(0.12),
                code_text, font_size=font_size, color=RGBColor(0xa5, 0xf3, 0xfc),
                word_wrap=True)

# ── Slide builders ─────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    # big logo-ish title centred
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
                "Unified Wi-Fi Mesh",
                font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
                "Controller DB Construct",
                font_size=30, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.5),
                "develop branch  ·  August 2025",
                font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    # bottom tag line
    add_textbox(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.4),
                "MariaDB persistence  ·  C++ multiple inheritance  ·  TR-181 bus integration",
                font_size=13, color=MUTED, align=PP_ALIGN.CENTER)
    return slide


def slide_block_diagram(prs):
    slide = blank_slide(prs)
    add_label(slide, "Block Diagram")
    add_title(slide, "Controller DB — System Block Diagram")
    add_divider(slide, Inches(1.2))

    def box(left, top, w, h, text, small="", color=CARD, border=ACCENT):
        b = slide.shapes.add_shape(1, left, top, w, h)
        b.fill.solid(); b.fill.fore_color.rgb = color
        b.line.color.rgb = border; b.line.width = Pt(1)
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = text
        run.font.bold = True; run.font.size = Pt(11); run.font.color.rgb = WHITE
        if small:
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = small
            r2.font.size = Pt(9); r2.font.color.rgb = MUTED
        return b

    def arrow_down(left, top, label="↓"):
        add_textbox(slide, left, top, Inches(1), Inches(0.28),
                    label, font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    def label_bar(top, text):
        add_textbox(slide, Inches(0.4), top, Inches(12.5), Inches(0.22),
                    text, font_size=8, bold=True,
                    color=RGBColor(0x64, 0x74, 0x8b))

    ROW1_TOP = Inches(1.3)
    # Row 1 — controller layer
    label_bar(ROW1_TOP - Inches(0.2), "CONTROLLER LAYER")
    box(Inches(0.4),  ROW1_TOP, Inches(2.2), Inches(0.75),
        "em_ctrl_t", "orchestration · state machines · 1905.1",
        color=RGBColor(0x1e,0x30,0x50), border=ACCENT)
    add_textbox(slide, Inches(2.65), ROW1_TOP + Inches(0.22), Inches(0.35), Inches(0.35),
                "↔", font_size=16, color=MUTED)
    box(Inches(3.0),  ROW1_TOP, Inches(5.2), Inches(0.75),
        "dm_easy_mesh_ctrl_t",
        "inherits all 11 list classes · owns m_db_client · init/load/update tables",
        color=RGBColor(0x1e,0x30,0x50), border=ACCENT)
    add_textbox(slide, Inches(8.25), ROW1_TOP + Inches(0.22), Inches(0.35), Inches(0.35),
                "↔", font_size=16, color=MUTED)
    box(Inches(8.6),  ROW1_TOP, Inches(4.1), Inches(0.75),
        "dm_easy_mesh_t [ ]",
        "per-network in-memory model (radios · BSS · STAs · policies)",
        color=RGBColor(0x1e,0x30,0x50), border=ACCENT)

    arrow_down(Inches(6.4), ROW1_TOP + Inches(0.78))

    ROW2_TOP = Inches(2.35)
    label_bar(ROW2_TOP - Inches(0.18), "DB PERSISTENCE CLASSES  (inherited by dm_easy_mesh_ctrl_t)")
    # 10 small boxes
    names = ["dm_network_list_t\nNetworkList",
             "dm_device_list_t\nDeviceList",
             "dm_radio_list_t\nRadioList",
             "dm_radio_cap_list_t\nCapabilities",
             "dm_bss_list_t\nBSSList",
             "dm_sta_list_t\nSTAList",
             "dm_policy_list_t\nPolicyList",
             "dm_network_ssid_list_t\nNetworkSSIDList",
             "dm_op_class_list_t\nOpClassList",
             "dm_scan_result_list_t\nScanResultList"]
    cell_w = Inches(1.25)
    for i, n in enumerate(names):
        parts = n.split("\n")
        b = slide.shapes.add_shape(1,
            Inches(0.4) + i * cell_w, ROW2_TOP, cell_w - Inches(0.05), Inches(0.65))
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0f,0x28,0x40)
        b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f); b.line.width = Pt(0.5)
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = parts[0]
        r.font.size = Pt(7.5); r.font.bold = True; r.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = parts[1] if len(parts) > 1 else ""
        r2.font.size = Pt(7); r2.font.color.rgb = MUTED

    arrow_down(Inches(6.4), ROW2_TOP + Inches(0.67), "↓  (all inherit)")

    ROW3_TOP = Inches(3.5)
    label_bar(ROW3_TOP - Inches(0.18), "DB PERSISTENCE BASE")
    box(Inches(0.4), ROW3_TOP, Inches(12.5), Inches(0.62),
        "db_easy_mesh_t  (DB persistence base)",
        "init_table · init_columns · sync_db · update_db · search_db · set/get_config",
        color=RGBColor(0x06,0x25,0x20), border=ACCENT2)

    arrow_down(Inches(6.4), ROW3_TOP + Inches(0.65), "↓  (calls via db_client_t value)")

    ROW4_TOP = Inches(4.52)
    label_bar(ROW4_TOP - Inches(0.18), "PERSISTENCE BACKEND")
    box(Inches(0.4), ROW4_TOP, Inches(12.5), Inches(0.62),
        "db_client_t",
        "concrete MariaDB/MySQL wrapper · MYSQL *m_con · libmariadb · on-device SQL",
        color=RGBColor(0x06,0x25,0x20), border=ACCENT2)

    arrow_down(Inches(6.4), ROW4_TOP + Inches(0.65))

    ROW5_TOP = Inches(5.5)
    label_bar(ROW5_TOP - Inches(0.18), "DATABASE")
    box(Inches(0.4), ROW5_TOP, Inches(12.5), Inches(0.55),
        "MariaDB",
        "10 SQL tables · persisted on-box · libmariadb",
        color=RGBColor(0x1a,0x1a,0x35), border=ACCENT)

    return slide


def slide_overview(prs):
    slide = blank_slide(prs)
    add_label(slide, "Overview")
    add_title(slide, "What is the Controller DB?")
    add_divider(slide, Inches(1.2))

    add_textbox(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.55),
                "The EasyMesh controller maintains a full in-memory replica of the mesh topology "
                "and persists it to a relational database. This 'data model' drives policy "
                "distribution, onboarding, and northbound TR-181 queries.",
                font_size=13, color=MUTED, word_wrap=True)

    cw = Inches(4.0)
    ch = Inches(1.5)
    tops = Inches(2.1)
    cards = [
        (ACCENT,  "Persistence Backend",
         "Single concrete db_client_t class wraps libmariadb (MYSQL *m_con). "
         "Owned by value inside dm_easy_mesh_ctrl_t — no factory, no polymorphism."),
        (ACCENT2, "DB Persistence Layer",
         "db_easy_mesh_t base class provides table init, CRUD, and JSON serialisation "
         "for every entity type."),
        (ACCENT3, "Domain Model",
         "dm_easy_mesh_ctrl_t inherits all list classes via C++ MI to form a single "
         "unified data model for the whole mesh network."),
    ]
    for i, (col, title, body) in enumerate(cards):
        add_card(slide, Inches(0.4) + i * (cw + Inches(0.2)), tops,
                 cw, ch, title, body, title_color=col, font_body=11)

    # Key design goal box
    add_card(slide, Inches(0.4), Inches(3.8), Inches(12.5), Inches(1.1),
             "Key Design Goal",
             "Keep persistence logic out of orchestration code — entity list classes "
             "receive db_client_t& by reference, so the controller core never embeds SQL directly.",
             title_color=ACCENT3, font_body=12)

    return slide


def slide_architecture(prs):
    slide = blank_slide(prs)
    add_label(slide, "Architecture")
    add_title(slide, "High-Level Layering")
    add_divider(slide, Inches(1.2))

    # Left: layer diagram
    LX = Inches(0.4)
    lw = Inches(5.8)

    def arch_box(top, h, text, small="", color=CARD, border=RGBColor(0x1e,0x3a,0x5f)):
        b = slide.shapes.add_shape(1, LX, top, lw, h)
        b.fill.solid(); b.fill.fore_color.rgb = color
        b.line.color.rgb = border; b.line.width = Pt(0.75)
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = text; r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = WHITE
        if small:
            p2 = tf.add_paragraph(); r2 = p2.add_run()
            r2.text = small; r2.font.size = Pt(9); r2.font.color.rgb = MUTED

    def arch_arrow(top):
        add_textbox(slide, LX + Inches(2.7), top, Inches(0.5), Inches(0.22),
                    "↓", font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

    y = Inches(1.35)
    arch_box(y, Inches(0.55), "em_ctrl_t",
             "Controller process — orchestration & message handling",
             color=RGBColor(0x0a,0x18,0x35), border=ACCENT)
    arch_arrow(y + Inches(0.56))
    y += Inches(0.78)
    arch_box(y, Inches(0.55), "dm_easy_mesh_ctrl_t",
             "Inherits all list classes — unified in-memory data model",
             color=RGBColor(0x1a,0x25,0x10), border=ACCENT3)
    arch_arrow(y + Inches(0.56))
    y += Inches(0.78)
    # mini list boxes (2 rows of 3)
    mnw = Inches(1.85); mnh = Inches(0.35)
    lists1 = ["dm_network_list_t", "dm_device_list_t", "dm_radio_list_t"]
    lists2 = ["dm_bss_list_t",     "dm_sta_list_t",    "dm_policy_list_t"]
    for j, n in enumerate(lists1):
        b = slide.shapes.add_shape(1, LX + j*(mnw+Inches(0.025)), y, mnw, mnh)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0f,0x1a,0x2a)
        b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f)
        tf = b.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = n; r.font.size = Pt(8); r.font.color.rgb = WHITE
    y += mnh + Inches(0.04)
    for j, n in enumerate(lists2):
        b = slide.shapes.add_shape(1, LX + j*(mnw+Inches(0.025)), y, mnw, mnh)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0f,0x1a,0x2a)
        b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f)
        tf = b.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = n; r.font.size = Pt(8); r.font.color.rgb = WHITE
    arch_arrow(y + mnh)
    y += mnh + Inches(0.22)
    arch_box(y, Inches(0.5), "db_easy_mesh_t",
             "DB persistence base — init / sync / update / search",
             color=RGBColor(0x06,0x20,0x18), border=ACCENT2)
    arch_arrow(y + Inches(0.51))
    y += Inches(0.72)
    arch_box(y, Inches(0.5), "db_client_t",
             "Concrete MariaDB wrapper · MYSQL *m_con · value member",
             color=RGBColor(0x06,0x20,0x18), border=ACCENT2)
    arch_arrow(y + Inches(0.51))
    y += Inches(0.72)
    arch_box(y, Inches(0.42), "MariaDB  —  10 SQL tables · libmariadb",
             color=RGBColor(0x1a,0x1a,0x35), border=ACCENT)

    # Right: cards
    RX = Inches(6.5)
    rw = Inches(6.45)
    ry = Inches(1.35)
    right_cards = [
        (ACCENT,  "Multiple Inheritance",
         "All list classes are inherited by dm_easy_mesh_ctrl_t via C++ MI, "
         "giving a single object that owns all tables simultaneously."),
        (ACCENT3, "Value Ownership",
         "db_client_t is a value member of dm_easy_mesh_ctrl_t. All list-class "
         "methods receive it by reference: update_db(m_db_client, …). "
         "No heap allocation; lifetime tied to the controller object."),
        (ACCENT2, "Per-Entity dm_*_t Objects",
         "Entity instances (dm_radio_t, dm_bss_t, dm_sta_t …) hold domain "
         "structs and are stored in dm_easy_mesh_t containers."),
    ]
    for title, body in [(c[1], c[2]) for c in right_cards]:
        col = right_cards[0][0] if title == right_cards[0][1] else \
              right_cards[1][0] if title == right_cards[1][1] else right_cards[2][0]
        add_card(slide, RX, ry, rw, Inches(1.55), title, body, title_color=col, font_body=11)
        ry += Inches(1.65)

    return slide


def slide_db_client(prs):
    slide = blank_slide(prs)
    add_label(slide, "Persistence Layer")
    add_title(slide, "db_client_t — Persistence Backend")
    add_divider(slide, Inches(1.2))

    code = (
        "// inc/db_client.h — concrete MariaDB class\n"
        "class db_client_t {\n"
        "  MYSQL *m_con;         // MariaDB connection\n"
        "  int connect(const char *path);\n"
        "public:\n"
        "  int   init(const char *path);   // format: \"user@password\"\n"
        "  void* execute(const char *query);\n"
        "  bool  next_result(void *ctx);\n"
        "  char* get_string(void *ctx, char *res, unsigned int col);\n"
        "  int   get_number(void *ctx, unsigned int col);\n"
        "  int   recreate_db();\n"
        "  db_client_t();\n"
        "  ~db_client_t();\n"
        "};\n\n"
        "// inc/dm_easy_mesh_ctrl.h — owned by value\n"
        "class dm_easy_mesh_ctrl_t : ... {\n"
        "  db_client_t m_db_client;  // value, not pointer\n"
        "};\n\n"
        "// Usage in table code (direct, no dereference)\n"
        "dm_network_list_t::load_table(m_db_client);\n"
        "dm_radio_list_t::update_db(m_db_client, op, data);"
    )
    add_code_box(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.6), code, font_size=9)

    add_card(slide, Inches(6.7), Inches(1.3), Inches(6.25), Inches(1.5),
             "Single Concrete Backend",
             "One class handles everything: connection pooling via MYSQL *m_con (libmariadb), "
             "SQL execution, and result iteration. MySQL header is conditionally included for "
             "OpenWRT/RDKB (<mysql/mysql.h>) vs Debian (<mariadb/mysql.h>).",
             title_color=ACCENT2)

    add_card(slide, Inches(6.7), Inches(3.0), Inches(6.25), Inches(1.35),
             "Value Ownership",
             "m_db_client is held as a value member of dm_easy_mesh_ctrl_t. All list-class "
             "calls pass it by reference: update_db(m_db_client, …). No heap allocation; "
             "lifetime tied to the controller object.",
             title_color=ACCENT)

    iter_code = (
        "void* ctx = m_db_client.execute(query);\n"
        "while (m_db_client.next_result(ctx)) {\n"
        "  m_db_client.get_string(ctx, buf, col);\n"
        "  m_db_client.get_number(ctx, col);\n"
        "}\n"
        "// next_result(false) frees ctx automatically"
    )
    add_textbox(slide, Inches(6.7), Inches(4.5), Inches(6.25), Inches(0.3),
                "Result Iteration Pattern", font_size=12, bold=True, color=ACCENT3)
    add_code_box(slide, Inches(6.7), Inches(4.82), Inches(6.25), Inches(1.05), iter_code)

    return slide


def slide_db_easy_mesh(prs):
    slide = blank_slide(prs)
    add_label(slide, "DB Persistence Layer")
    add_title(slide, "db_easy_mesh_t — Table Base Class")
    add_divider(slide, Inches(1.2))

    code = (
        "class db_easy_mesh_t {\n"
        "public:\n"
        "  db_table_name_t  m_table_name;   // e.g. \"RadioList\"\n"
        "  unsigned int     m_num_cols;\n"
        "  db_column_t      m_columns[EM_MAX_COLS];\n\n"
        "  // Schema bootstrap\n"
        "  virtual int  init_table(db_client_t&);\n"
        "  virtual int  init_columns(db_client_t&);\n\n"
        "  // CRUD\n"
        "  virtual int  sync_db(db_client_t&, void *ctx);\n"
        "  virtual int  update_db(db_client_t&, dm_orch_type_t op, void *data);\n"
        "  virtual bool search_db(db_client_t&, void *ctx, void *key);\n\n"
        "  // JSON serialisation\n"
        "  virtual int    set_config(db_client_t&, const cJSON*, void *parent_id);\n"
        "  virtual cJSON* get_config(db_client_t&, void *parent_id);\n"
        "};"
    )
    add_code_box(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(4.8), code)

    add_card(slide, Inches(6.7), Inches(1.3), Inches(6.25), Inches(1.2),
             "Column Descriptor",
             "Each column is described by db_column_t: name (char[64]), type "
             "(db_column_type_varchar/int/…), size, and a value buffer. "
             "init_columns() fills this array; init_table() builds the CREATE TABLE SQL from it.",
             title_color=ACCENT)
    add_card(slide, Inches(6.7), Inches(2.65), Inches(6.25), Inches(1.1),
             "Orchestration Type",
             "dm_orch_type_t passed to update_db() selects INSERT (db_create), "
             "UPDATE (db_update), or DELETE (db_delete). Each list class handles "
             "all three cases in one override.",
             title_color=ACCENT3)
    add_card(slide, Inches(6.7), Inches(3.9), Inches(6.25), Inches(1.1),
             "Table Bootstrap",
             "Each list class calls init() which runs CREATE TABLE IF NOT EXISTS via "
             "db_client_t::execute(). Called before m_db_client.init() so the schema "
             "exists before the MariaDB connection is established.",
             title_color=ACCENT2)

    return slide


def slide_entity_list(prs):
    slide = blank_slide(prs)
    add_label(slide, "Design Pattern")
    add_title(slide, "Entity + List Dual-Class Pattern")
    add_divider(slide, Inches(1.2))

    add_textbox(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.45),
                "Every domain entity uses a two-class split: a pure-data entity class "
                "and a DB-aware list class that inherits it.",
                font_size=13, color=MUTED, word_wrap=True)

    code = (
        "// Pure domain — shared with agent\n"
        "class dm_radio_t {\n"
        "  em_radio_info_t m_radio_info;\n"
        "  // ... accessors, capability decode ...\n"
        "};\n\n"
        "// DB-aware — controller only\n"
        "class dm_radio_list_t : public dm_radio_t,\n"
        "                        public db_easy_mesh_t {\n"
        "  int  sync_db(db_client_t&, void*)  override;\n"
        "  int  update_db(db_client_t&, dm_orch_type_t, void*) override;\n"
        "  bool search_db(db_client_t&, void*, void*) override;\n"
        "  int  init_columns(db_client_t&)    override;\n"
        "};\n\n"
        "// dm_easy_mesh_ctrl_t inherits ALL list classes:\n"
        "class dm_easy_mesh_ctrl_t\n"
        "  : public dm_network_list_t,\n"
        "    public dm_device_list_t,\n"
        "    public dm_radio_list_t,\n"
        "    // ... 8 more list classes ...\n"
        "    public dm_scan_result_list_t {\n"
        "  db_client_t m_db_client;  // value\n"
        "};"
    )
    add_code_box(slide, Inches(0.4), Inches(1.9), Inches(5.8), Inches(5.1), code, font_size=9)

    RX = Inches(6.5)
    add_card(slide, RX, Inches(1.9), Inches(6.45), Inches(1.55),
             "Benefit of the Split",
             "Entity classes can be used on the agent (no DB) with zero overhead.\n"
             "List classes add table schema and SQL logic — controller-only binary.\n"
             "Agent-side code depends only on dm_radio_t, never on dm_radio_list_t.",
             title_color=ACCENT2)

    add_textbox(slide, RX, Inches(3.6), Inches(6.45), Inches(0.3),
                "Entities That Follow This Pattern", font_size=12, bold=True, color=ACCENT)
    pairs = [
        ("dm_network_t",    "dm_network_list_t"),
        ("dm_device_t",     "dm_device_list_t"),
        ("dm_radio_t",      "dm_radio_list_t"),
        ("dm_radio_cap_t",  "dm_radio_cap_list_t"),
        ("dm_bss_t",        "dm_bss_list_t"),
        ("dm_sta_t",        "dm_sta_list_t"),
        ("dm_policy_t",     "dm_policy_list_t"),
        ("dm_network_ssid_t","dm_network_ssid_list_t"),
        ("dm_op_class_t",   "dm_op_class_list_t"),
        ("dm_scan_result_t","dm_scan_result_list_t"),
    ]
    bw = Inches(3.1)
    for i, (entity, lst) in enumerate(pairs):
        row = i // 2
        col = i % 2
        bx = RX + col * (bw + Inches(0.12))
        by = Inches(4.0) + row * Inches(0.42)
        b = slide.shapes.add_shape(1, bx, by, bw, Inches(0.36))
        b.fill.solid(); b.fill.fore_color.rgb = CARD
        b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f)
        tf = b.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = f"{entity}  →  {lst}"
        r.font.size = Pt(9); r.font.color.rgb = WHITE

    return slide


def slide_schema_overview(prs):
    slide = blank_slide(prs)
    add_label(slide, "Schema")
    add_title(slide, "Database Tables — At a Glance")
    add_divider(slide, Inches(1.2))

    rows = [
        ("NetworkList",        "dm_network_list_t",        "ID (char 64)",  "ControllerID, ColocatedAgentID, Media"),
        ("DeviceList",         "dm_device_list_t",         "ID (char 128)", "DeviceID, Profile, Manufacturer, SoftwareVersion, CountryCode"),
        ("RadioList",          "dm_radio_list_t",          "ID (char 128)", "RadioID, Band, CF1/CF2, Utilization, Noise, NumberOfBSS"),
        ("Capabilities",       "dm_radio_cap_list_t",      "ID (char 17)",  "HTCapabilities, VHTCapabilities, HECapabilities, EHTCapabilities"),
        ("BSSList",            "dm_bss_list_t",            "ID (char 128)", "BSSID, RUID, SSID, FronthaulUse, BackhaulUse, VlanID"),
        ("STAList",            "dm_sta_list_t",            "MACAddress",    "BSSID, RUID, SignalStrength, RCPI, BytesSent/Received, Associated"),
        ("PolicyList",         "dm_policy_list_t",         "ID (char 64)",  "PolicyType, APMetricsInterval, SteeringPolicyType, STARCPIThresold"),
        ("NetworkSSIDList",    "dm_network_ssid_list_t",   "ID (char 64)",  "SSID, PassPhrase, Band, AKMsAllowed, HaulType, VLANID"),
        ("OperatingClassList", "dm_op_class_list_t",       "ID (char 32)",  "Class, Channel, ChannelList, TxPower, MaxTxPower"),
        ("ScanResultList",     "dm_scan_result_list_t",    "ID (char 128)", "BSSID, SSID, SignalStrength, Utilization, STACount, ScanType"),
    ]

    # header
    HDR_H = Inches(0.32)
    TBL_TOP = Inches(1.32)
    COL_W = [Inches(1.65), Inches(2.5), Inches(1.6), Inches(6.55)]
    COL_X = [Inches(0.4), Inches(2.08), Inches(4.62), Inches(6.25)]
    hdrs = ["Table Name", "List Class", "Primary Key", "Key Columns"]
    for i, hdr in enumerate(hdrs):
        b = slide.shapes.add_shape(1, COL_X[i], TBL_TOP, COL_W[i], HDR_H)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0a,0x20,0x3a)
        b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f)
        tf = b.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = hdr; r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = ACCENT

    ROW_H = Inches(0.44)
    for ri, (tbl, cls, pk, cols) in enumerate(rows):
        row_top = TBL_TOP + HDR_H + ri * ROW_H
        bg = RGBColor(0x0f,0x1a,0x2a) if ri % 2 == 0 else CARD
        values = [tbl, cls, pk, cols]
        for ci, val in enumerate(values):
            b = slide.shapes.add_shape(1, COL_X[ci], row_top, COL_W[ci], ROW_H)
            b.fill.solid(); b.fill.fore_color.rgb = bg
            b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f); b.line.width = Pt(0.25)
            tf = b.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = val
            r.font.size = Pt(9)
            r.font.color.rgb = ACCENT2 if ci == 0 else WHITE

    return slide


def slide_core_tables(prs):
    slide = blank_slide(prs)
    add_label(slide, "Schema Deep-Dive")
    add_title(slide, "Core RF Topology Tables")
    add_divider(slide, Inches(1.2))

    radio_cols = (
        "• RadioID — RUID (MAC)\n"
        "• Band — 2.4 / 5 / 6 GHz\n"
        "• CF1, CF2 — centre frequencies\n"
        "• Utilization, Noise\n"
        "• NumberOfBSS\n"
        "• SteeringPolicy\n"
        "• RCPISteeringThreshold\n"
        "• STAReportingRCPIThreshold\n"
        "• TrafficSeparationCombined*\n"
        "• ChipsetVendor"
    )
    bss_cols = (
        "• BSSID — BSS MAC\n"
        "• RUID → RadioList FK\n"
        "• SSID, Enabled\n"
        "• FronthaulUse, BackhaulUse\n"
        "• FronthaulAKMsAllowed\n"
        "• BackhaulAKMsAllowed\n"
        "• Profile1bSTAsDisallowed\n"
        "• Profile2bSTAsDisallowed\n"
        "• VlanID\n"
        "• EstServiceParameters{BE/BK/VI/VO}"
    )
    sta_cols = (
        "• MACAddress — STA MAC (PK)\n"
        "• BSSID → BSSList FK\n"
        "• RUID → RadioList FK\n"
        "• Associated flag\n"
        "• SignalStrength, RCPI\n"
        "• LastDataUplinkRate / DownlinkRate\n"
        "• PacketsSent, BytesSent\n"
        "• PacketsReceived, BytesReceived\n"
        "• RetransCount\n"
        "• FrameBodyLength + FrameBody\n"
        "  (re)assoc request, hex-encoded"
    )
    cw = Inches(4.1); ch = Inches(4.0)
    add_card(slide, Inches(0.4),  Inches(1.32), cw, ch,
             "RadioList  —  one row per radio per device", radio_cols, title_color=ACCENT)
    add_card(slide, Inches(4.6),  Inches(1.32), cw, ch,
             "BSSList  —  one row per BSS per radio", bss_cols, title_color=ACCENT2)
    add_card(slide, Inches(8.8),  Inches(1.32), cw, ch,
             "STAList  —  one row per associated STA", sta_cols, title_color=ACCENT3)

    # relationship note
    add_card(slide, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.8),
             "Relationship Chain",
             "NetworkList → DeviceList → RadioList → BSSList → STAList  "
             "(hierarchy via composite ID keys; no explicit FK constraints in the schema)",
             title_color=MUTED, font_body=12)
    return slide


def slide_policy_tables(prs):
    slide = blank_slide(prs)
    add_label(slide, "Schema Deep-Dive")
    add_title(slide, "Policy & Network Configuration Tables")
    add_divider(slide, Inches(1.2))

    policy_rows = [
        ("PolicyType",              "steering_local / btm / param / ap_metrics / radio_metrics …"),
        ("APMetricsInterval",       "AP metrics collection interval (seconds)"),
        ("STAList",                 "Serialised MAC list (TEXT)"),
        ("SteeringPolicyType",      "mandate / disallow / allow"),
        ("STARCPIThresold",         "RCPI threshold for STA steering"),
        ("STARCPIHysteresis",       "Hysteresis margin override"),
        ("APUtilThreshold",         "Channel utilisation threshold"),
        ("IndependentScanRep",      "Channel scan report flag"),
        ("Profile_1/2_Disallowed",  "Backhaul BSS profile flags"),
        ("PrimaryVlanId",           "Default 802.1Q VLAN"),
        ("ManagedClientMarker",     "Vendor OUI marker string"),
    ]

    add_textbox(slide, Inches(0.4), Inches(1.32), Inches(6.5), Inches(0.3),
                "PolicyList", font_size=13, bold=True, color=ACCENT4)
    HDR_H = Inches(0.28); ROW_H = Inches(0.38); TBL_TOP = Inches(1.65)
    for i, hdr in enumerate(["Column", "Purpose"]):
        b = slide.shapes.add_shape(1, Inches(0.4) + i*Inches(3.2), TBL_TOP,
                                   Inches(3.1), HDR_H)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0a,0x20,0x3a)
        b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f)
        tf = b.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = hdr; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = ACCENT
    for ri, (col, purpose) in enumerate(policy_rows):
        rt = TBL_TOP + HDR_H + ri * ROW_H
        bg = RGBColor(0x0f,0x1a,0x2a) if ri % 2 == 0 else CARD
        for ci, val in enumerate([col, purpose]):
            b = slide.shapes.add_shape(1, Inches(0.4)+ci*Inches(3.2), rt, Inches(3.1), ROW_H)
            b.fill.solid(); b.fill.fore_color.rgb = bg
            b.line.color.rgb = RGBColor(0x1e,0x3a,0x5f); b.line.width = Pt(0.25)
            tf = b.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = val; r.font.size = Pt(8.5)
            r.font.color.rgb = ACCENT2 if ci == 0 else WHITE

    RX = Inches(6.9)
    add_card(slide, RX, Inches(1.32), Inches(6.0), Inches(2.0),
             "NetworkSSIDList",
             "• SSID, PassPhrase\n"
             "• Band (2.4 / 5 / 6 GHz bitmask)\n"
             "• AKMsAllowed — WPA2/WPA3/OWE etc.\n"
             "• HaulType — fronthaul / backhaul\n"
             "• MFPConfig — 802.11w setting\n"
             "• AuthType, SuiteSelector\n"
             "• VLANID — traffic separation",
             title_color=ACCENT3)
    add_card(slide, RX, Inches(3.45), Inches(6.0), Inches(2.0),
             "OperatingClassList",
             "• Class — 802.11 op class number\n"
             "• Channel — current channel\n"
             "• ChannelList — usable channels (varchar 256)\n"
             "• ChannelPreferenceList — scored preferences\n"
             "• TxPower, MaxTxPower\n"
             "• Minutes, Seconds, Countdown — CAC timer",
             title_color=ACCENT)
    return slide


def slide_init_lifecycle(prs):
    slide = blank_slide(prs)
    add_label(slide, "Lifecycle")
    add_title(slide, "Controller DB Initialization Flow")
    add_divider(slide, Inches(1.2))

    # flow strip
    steps = ["em_ctrl_t::init()", "dm_easy_mesh_ctrl_t::init(path,mgr)",
             "init_tables()", "m_db_client.init(path)",
             "tr_181_t::init(this)", "load_tables()"]
    sw = Inches(2.0); sx = Inches(0.4); sy = Inches(1.3); sh = Inches(0.42)
    for i, step in enumerate(steps):
        b = slide.shapes.add_shape(1, sx + i*(sw+Inches(0.18)), sy, sw, sh)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0a,0x18,0x35)
        b.line.color.rgb = ACCENT; b.line.width = Pt(0.75)
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = step; r.font.size = Pt(8.5); r.font.color.rgb = WHITE
        if i < len(steps)-1:
            add_textbox(slide, sx + i*(sw+Inches(0.18)) + sw, sy + Inches(0.06),
                        Inches(0.18), sh, "→", font_size=13, color=MUTED)

    init_code = (
        "void init_tables() {\n"
        "  dm_network_list_t::    init();  // CREATE TABLE IF NOT EXISTS\n"
        "  dm_device_list_t::     init();\n"
        "  dm_radio_list_t::      init();\n"
        "  dm_network_ssid_list_t::init();\n"
        "  dm_radio_cap_list_t::  init();\n"
        "  dm_op_class_list_t::   init();\n"
        "  dm_bss_list_t::        init();\n"
        "  dm_sta_list_t::        init();\n"
        "  dm_policy_list_t::     init();\n"
        "  dm_scan_result_list_t::init();\n"
        "}\n"
        "// Called BEFORE m_db_client.init() so tables\n"
        "// exist before the MariaDB connection opens."
    )
    add_textbox(slide, Inches(0.4), Inches(2.0), Inches(6.0), Inches(0.3),
                "init_tables() — Schema Bootstrap", font_size=12, bold=True, color=ACCENT3)
    add_code_box(slide, Inches(0.4), Inches(2.33), Inches(6.0), Inches(3.7), init_code)

    load_code = (
        "int load_tables() {\n"
        "  // SELECT * → sync_db() per row for each table\n"
        "  if (dm_network_list_t::load_table(m_db_client) != 0)\n"
        "    return db_cfg_type_network_list_update;\n"
        "  if (dm_device_list_t::load_table(m_db_client) != 0)\n"
        "    return db_cfg_type_device_list_update;\n"
        "  // ... radio, bss, sta, policy, op_class ...\n"
        "  if (is_table_empty(m_db_client)) return -1;\n"
        "  set_initialized();\n"
        "  return 0;\n"
        "}\n\n"
        "// Empty DB → run seed script → reload\n"
        "if (rc == -1) {\n"
        "  system(\"/usr/ccsp/EasyMesh/setup_mysql_db_post.sh\");\n"
        "  rc = load_tables();\n"
        "}"
    )
    add_textbox(slide, Inches(6.7), Inches(2.0), Inches(6.25), Inches(0.3),
                "load_tables() — Populate from DB", font_size=12, bold=True, color=ACCENT2)
    add_code_box(slide, Inches(6.7), Inches(2.33), Inches(6.25), Inches(3.0), load_code)
    add_card(slide, Inches(6.7), Inches(5.45), Inches(6.25), Inches(0.85),
             "First-Boot Seed",
             "If tables are empty on first boot the controller auto-runs the seed script, "
             "then calls load_tables() a second time. If still fails, init() returns -1 and the process exits.",
             title_color=ACCENT3, font_body=10)
    return slide


def slide_update_flow(prs):
    slide = blank_slide(prs)
    add_label(slide, "Lifecycle")
    add_title(slide, "Runtime Update Flow — How Changes Reach the DB")
    add_divider(slide, Inches(1.2))

    steps = ["1905 Message\n/ Event Bus", "em_*_t\nhandler",
             "Update dm_easy_mesh_t", "Set dirty flag\ndb_cfg_type_t",
             "update_tables(dm)", "update_db\n(client, op, data)", "SQL INSERT\n/UPDATE/DELETE"]
    sw = Inches(1.72); sx = Inches(0.4); sy = Inches(1.3); sh = Inches(0.65)
    for i, step in enumerate(steps):
        b = slide.shapes.add_shape(1, sx + i*(sw+Inches(0.1)), sy, sw, sh)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x0a,0x18,0x35)
        b.line.color.rgb = ACCENT; b.line.width = Pt(0.75)
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = step; r.font.size = Pt(8); r.font.color.rgb = WHITE
        if i < len(steps)-1:
            add_textbox(slide, sx + i*(sw+Inches(0.1)) + sw, sy + Inches(0.16),
                        Inches(0.1), sh, "→", font_size=12, color=MUTED)

    code = (
        "int update_tables(dm_easy_mesh_t *dm) {\n"
        "  // Network\n"
        "  dm_network_list_t::update_db(\n"
        "    *m_db_client, dm_orch_type_db_update,\n"
        "    dm->get_network_info());\n\n"
        "  // Device\n"
        "  dm_device_list_t::update_db(\n"
        "    *m_db_client, dm_orch_type_db_update,\n"
        "    dm->get_device_info());\n\n"
        "  // Radios (per radio)\n"
        "  for (int i = 0; i < dm->get_num_radios(); i++)\n"
        "    dm_radio_list_t::update_db(\n"
        "      *m_db_client, dm_orch_type_db_update,\n"
        "      dm->get_radio(i)->get_radio_info());\n\n"
        "  // BSS, STA, OpClass, ScanResult …\n"
        "}"
    )
    add_textbox(slide, Inches(0.4), Inches(2.1), Inches(6.0), Inches(0.3),
                "update_tables() Orchestration", font_size=12, bold=True, color=ACCENT3)
    add_code_box(slide, Inches(0.4), Inches(2.43), Inches(6.0), Inches(4.2), code)

    dirty_code = (
        "db_cfg_type_network_list_update\n"
        "db_cfg_type_device_list_update\n"
        "db_cfg_type_radio_list_update\n"
        "db_cfg_type_bss_list_update\n"
        "db_cfg_type_sta_list_update\n"
        "db_cfg_type_policy_list_update\n"
        "// … one per table"
    )
    add_card(slide, Inches(6.7), Inches(2.1), Inches(6.25), Inches(2.0),
             "db_cfg_type_t — Dirty Flags",
             "Each table has a corresponding flag bit. The orchestration layer tests these "
             "to skip tables that have not changed, avoiding unnecessary SQL traffic.",
             title_color=ACCENT)
    add_code_box(slide, Inches(6.7), Inches(3.55), Inches(6.25), Inches(1.0), dirty_code, font_size=9)

    add_card(slide, Inches(6.7), Inches(4.7), Inches(6.25), Inches(1.0),
             "delete vs update vs create",
             "Caller selects dm_orch_type_db_delete, dm_orch_type_db_update, or "
             "dm_orch_type_db_create. Each list class translates this to the appropriate SQL.",
             title_color=ACCENT2, font_body=11)
    return slide


def slide_in_memory(prs):
    slide = blank_slide(prs)
    add_label(slide, "In-Memory Model")
    add_title(slide, "dm_easy_mesh_t — Per-Network Container")
    add_divider(slide, Inches(1.2))

    add_textbox(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.45),
                "While dm_easy_mesh_ctrl_t owns the DB tables, the live in-memory state for each "
                "active mesh network lives in dm_easy_mesh_t instances (one per managed network).",
                font_size=13, color=MUTED, word_wrap=True)

    code = (
        "class dm_easy_mesh_t {\n"
        "  // Scalar state\n"
        "  em_network_info_t    m_network;\n"
        "  em_device_info_t     m_device;\n\n"
        "  // Arrays\n"
        "  dm_radio_t           m_radio[EM_MAX_RADIOS];\n"
        "  dm_bss_t             m_bss[EM_MAX_BSS_PER_RADIO * EM_MAX_RADIOS];\n"
        "  dm_network_ssid_t    m_network_ssid[EM_MAX_NETWORK_SSID];\n"
        "  dm_op_class_t        m_op_class[EM_MAX_OP_CLASS];\n"
        "  dm_policy_t          m_policy[EM_MAX_POLICIES];\n\n"
        "  // STAs — hash map for O(1) lookup\n"
        "  unordered_map<mac_address_t, dm_sta_t*> m_sta_map;\n"
        "  map<mac_address_t, dm_sta_t*>           m_sta_assoc_map;\n"
        "  map<mac_address_t, dm_sta_t*>           m_sta_dassoc_map;\n\n"
        "  // Helpers\n"
        "  em_radio_info_t* get_radio_info(int idx);\n"
        "  unsigned int     get_num_radios();\n"
        "  unsigned int     get_num_policy();\n"
        "  bool has_policy_type(em_policy_id_type_t t);\n"
        "};"
    )
    add_code_box(slide, Inches(0.4), Inches(1.9), Inches(6.0), Inches(4.75), code)

    add_card(slide, Inches(6.7), Inches(1.9), Inches(6.25), Inches(1.5),
             "Controller vs Agent Usage",
             "dm_easy_mesh_t is shared code. The controller uses it as the in-memory "
             "representation that the DB is sync'd to/from. The agent uses it as a local "
             "cache of topology received from the controller, with no DB involvement.",
             title_color=ACCENT)
    add_card(slide, Inches(6.7), Inches(3.55), Inches(6.25), Inches(1.35),
             "STA Tracking Maps",
             "• m_sta_map — all known STAs (O(1) MAC lookup)\n"
             "• m_sta_assoc_map — currently associated STAs\n"
             "• m_sta_dassoc_map — recently disassociated STAs (history)",
             title_color=ACCENT3)
    add_card(slide, Inches(6.7), Inches(5.05), Inches(6.25), Inches(1.0),
             "Policy Helper",
             "has_policy_type(t) lets send_policy_cfg_request_msg() skip TLV construction "
             "for policy types not present — enabling incremental policy updates.",
             title_color=ACCENT2, font_body=11)
    return slide


def slide_integration(prs):
    slide = blank_slide(prs)
    add_label(slide, "Integration")
    add_title(slide, "JSON Serialisation & TR-181 Bus Integration")
    add_divider(slide, Inches(1.2))

    flow_code = (
        "// Inbound — JSON → DB\n"
        "set_config(db_client, json_obj, parent_id)\n"
        "  → entity::decode(json_obj, parent_id)\n"
        "  → update_db(client, dm_orch_type_db_create, data)\n\n"
        "// Outbound — DB → JSON\n"
        "get_config(db_client, parent_id)\n"
        "  → entity::encode()\n"
        "  → returns cJSON* tree"
    )
    add_textbox(slide, Inches(0.4), Inches(1.32), Inches(6.0), Inches(0.3),
                "JSON In/Out via set_config / get_config", font_size=12, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.4), Inches(1.65), Inches(6.0), Inches(0.5),
                "Every entity implements decode(cJSON*, void*) and encode(). "
                "The controller ingests JSON commands from northbound (CLI/WebUI) and maps them to DB rows.",
                font_size=11, color=MUTED, word_wrap=True)
    add_code_box(slide, Inches(0.4), Inches(2.2), Inches(6.0), Inches(1.8), flow_code)

    # command flow
    cmd_steps = ["JSON cmd", "em_cmd_t", "dm analysis", "set_config", "update_db", "1905 policy msg"]
    csw = Inches(1.85); csx = Inches(0.4); csy = Inches(4.15); csh = Inches(0.4)
    for i, step in enumerate(cmd_steps):
        b = slide.shapes.add_shape(1, csx + i*(csw+Inches(0.12)), csy, csw, csh)
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x06,0x20,0x18)
        b.line.color.rgb = ACCENT2; b.line.width = Pt(0.5)
        tf = b.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = step; r.font.size = Pt(9); r.font.color.rgb = WHITE
        if i < len(cmd_steps)-1:
            add_textbox(slide, csx + i*(csw+Inches(0.12)) + csw, csy + Inches(0.06),
                        Inches(0.12), csh, "→", font_size=11, color=MUTED)

    # TR-181 handlers card
    add_card(slide, Inches(6.7), Inches(1.32), Inches(6.25), Inches(2.9),
             "TR-181 Bus Handlers",
             "dm_easy_mesh_ctrl_t registers bus get/set callbacks for every TR-181 path:\n\n"
             "• network_get / device_get\n"
             "• radio_get / radio_tget\n"
             "• bss_get / bss_tget\n"
             "• ssid_get / rcaps_get\n"
             "• sta_get / policy_get\n"
             "• X_AIRTIES_UnassociatedStaLinkMetricsQuery()  ← new in develop\n\n"
             "Handlers read from in-memory model (not SQL at query time) to keep latency low.",
             title_color=ACCENT)

    # bottom 3 cards
    bw = Inches(4.0)
    bottom_cards = [
        (ACCENT3, "db_cfg_type_t Flags",
         "Set when JSON changes are parsed. Drive which tables get synced "
         "and which 1905.1 messages are constructed."),
        (ACCENT,  "Multi-Network",
         "dm_easy_mesh_list_t holds a list of dm_easy_mesh_t* — one per managed "
         "EasyMesh network ID. The controller handles them concurrently."),
        (ACCENT2, "Thread-Safety Note",
         "The DB client interface explicitly documents that implementations are "
         "not required to be thread-safe. The orchestration loop is single-threaded."),
    ]
    for i, (col, title, body) in enumerate(bottom_cards):
        add_card(slide, Inches(0.4) + i*(bw+Inches(0.2)), Inches(4.75), bw, Inches(1.5),
                 title, body, title_color=col, font_body=10)
    return slide


def slide_summary(prs):
    slide = blank_slide(prs)
    add_label(slide, "Summary")
    add_title(slide, "Key Takeaways")
    add_divider(slide, Inches(1.2))

    left_cards = [
        (ACCENT,  "1. Single Concrete DB Backend",
         "db_client_t is a concrete MariaDB wrapper with MYSQL *m_con. Owned by value in "
         "dm_easy_mesh_ctrl_t. All 10 list classes call it by reference. MySQL header is "
         "conditionally included for OpenWRT/RDKB vs Debian."),
        (ACCENT2, "2. Entity + List Split",
         "Pure-domain dm_*_t classes are shared with the agent. DB-aware dm_*_list_t "
         "subclasses are controller-only. Zero persistence overhead on the agent side."),
        (ACCENT3, "3. Unified Controller Object",
         "dm_easy_mesh_ctrl_t inherits all 11 list classes via C++ MI — a single object "
         "that can query, update, or delete any table without any additional glue code."),
    ]
    right_cards = [
        (ACCENT4, "4. Dirty-Flag Driven Sync",
         "db_cfg_type_t bitmask tracks exactly which tables changed. update_tables() "
         "skips clean tables — minimising SQL round-trips during high-frequency topology updates."),
        (ACCENT,  "5. 10 Tables, Clear Hierarchy",
         "NetworkList → DeviceList → RadioList → Capabilities → BSSList → STAList\n"
         "NetworkSSIDList · PolicyList · OperatingClassList · ScanResultList"),
        (ACCENT2, "6. JSON ↔ DB ↔ 1905.1 Pipeline",
         "Northbound JSON commands flow through set_config → update_db → dirty flags "
         "→ 1905.1 TLV construction. Southbound telemetry flows the other way via sync_db "
         "and TR-181 bus handlers."),
    ]
    cw = Inches(6.1); ch = Inches(1.35)
    for i, (col, title, body) in enumerate(left_cards):
        add_card(slide, Inches(0.4), Inches(1.3) + i*(ch+Inches(0.12)),
                 cw, ch, title, body, title_color=col, font_body=10)
    for i, (col, title, body) in enumerate(right_cards):
        add_card(slide, Inches(6.8), Inches(1.3) + i*(ch+Inches(0.12)),
                 cw, ch, title, body, title_color=col, font_body=10)

    # Recent additions
    add_card(slide, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.75),
             "7. Recent Additions (develop branch)",
             "RDKBWIFI-555: dm_sta_t::get_assoc_frame_ie_offset() probes assoc/reassoc/IEs-only "
             "layouts to find where IEs start. Fixes misparse of capabilities for reassociating STAs.\n"
             "RDKBWIFI-529: decode_sta_capability() now called inside TR-181 sta_get/sta_tget before "
             "derived fields are read. Fixes alternating empty HTCapabilities under polling.\n"
             "RDKBWIFI-408: X_AIRTIES_UnassociatedStaLinkMetricsQuery() TR-181 method added. "
             "Accepts OpClass + channel/STA lists, builds UnassocSTAQuery JSON, enqueues as em_cmd_t.",
             title_color=ACCENT4, font_body=10)
    return slide


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_title(prs)
    slide_block_diagram(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_db_client(prs)
    slide_db_easy_mesh(prs)
    slide_entity_list(prs)
    slide_schema_overview(prs)
    slide_core_tables(prs)
    slide_policy_tables(prs)
    slide_init_lifecycle(prs)
    slide_update_flow(prs)
    slide_in_memory(prs)
    slide_integration(prs)
    slide_summary(prs)

    out = "/root/build/mar_31/openwrt/package/easymesh/unified-wifi-mesh/docs/db_construct_slides.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")

if __name__ == "__main__":
    main()
